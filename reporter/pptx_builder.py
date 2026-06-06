"""
PPTX Builder Module
使用 python-pptx 生成专业地质勘探演示文稿。
深色主题 + 金色强调色，图表为主，16:9 宽屏。
"""

import io
import re
import tempfile
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt as Pt2
from pptx.oxml.ns import qn

from .categories import SearchResult, get_all_categories
from .geocoder import LocationContext

# 无公开数据的占位特征——命中则该数据点不展示（删除占位句，不写"暂缺/来源"）
_NO_DATA_MARKERS = ("暂缺", "暂无", "无数据", "无相关", "未获取", "未检索", "n/a", "N/A", "待补充", "缺失")


def _is_no_data(value) -> bool:
    if not value or not str(value).strip():
        return True
    return any(m in str(value) for m in _NO_DATA_MARKERS)


def _valid_dps(result):
    """返回有公开数据的数据点（过滤占位行）。"""
    return [dp for dp in getattr(result, "data_points", []) if not _is_no_data(dp.value)]


class PptxBuilder:
    """生成专业地质勘探 PPT 演示文稿"""

    # 配色方案
    COLOR_BG = RGBColor(0x1B, 0x2A, 0x4A)         # 深蓝灰背景
    COLOR_BG_LIGHT = RGBColor(0x24, 0x3B, 0x63)    # 稍浅背景（卡片）
    COLOR_ACCENT = RGBColor(0xD4, 0xA8, 0x43)      # 金色强调
    COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    COLOR_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)       # 浅灰文本
    COLOR_TABLE_HEADER = RGBColor(0x2C, 0x3E, 0x6B) # 表头深色
    COLOR_TABLE_ROW1 = RGBColor(0x1F, 0x30, 0x55)  # 表格奇数行
    COLOR_TABLE_ROW2 = RGBColor(0x17, 0x25, 0x42)  # 表格偶数行

    # 字体
    FONT_TITLE = "微软雅黑"
    FONT_BODY = "宋体"

    # 尺寸（16:9）
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def __init__(self, output_dir: str = "./reports"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def build_pptx(self, location: LocationContext, search_results: Dict[str, SearchResult],
                   output_name: str = None, mineral_type: str = "",
                   target_figure=None, confidence: dict = None) -> str:
        """
        生成完整的 PPT 演示文稿。

        Returns
        -------
        str
            生成的 .pptx 文件路径
        """
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT

        # 封面
        self._add_cover_slide(prs, location, mineral_type)

        # 研究区概况
        self._add_overview_slide(prs, location, mineral_type)

        # 每个数据类别一页
        for category in get_all_categories():
            result = search_results.get(category.id)
            if result and not result.error:
                self._add_category_slide(prs, category.chapter_title, result, mineral_type, category.id)
            # 子系统图件单独成页（蚀变/物探/深部预测等）
            _figs = getattr(result, "figures", None) or [] if result else []
            if _figs:
                self._add_figures_slide(prs, category.chapter_title, _figs)

        # 靶区推荐页
        self._add_target_slide(prs, location, target_figure, mineral_type)

        # 综合置信评价页
        self._add_confidence_slide(prs, confidence, mineral_type)

        # 综合发现汇总
        self._add_summary_slide(prs, search_results, mineral_type)

        # 保存
        if not output_name:
            timestamp = datetime.now().strftime("%Y%m%d")
            output_name = f"{location.area_name}_{timestamp}"

        output_path = self.output_dir / f"{output_name}.pptx"
        prs.save(str(output_path))
        return str(output_path)

    # ------------------------------------------------------------------
    # 内部方法
    # ------------------------------------------------------------------

    def _set_slide_bg(self, slide, color=None):
        """设置幻灯片背景色"""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color or self.COLOR_BG

    def _add_textbox(self, slide, left, top, width, height, text, font_size=18,
                     font_name=None, color=None, bold=False, alignment=PP_ALIGN.LEFT):
        """添加文本框"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = font_name or self.FONT_BODY
        p.font.color.rgb = color or self.COLOR_WHITE
        p.font.bold = bold
        p.alignment = alignment
        return txBox

    def _set_paragraph_with_numeric_bold(self, p, text: str, base_font_size: int,
                                          font_name: str, color):
        """将段落文本中的数字拆成独立 run，数字部分加大加粗显示。"""
        from pptx.util import Pt
        from pptx.oxml.ns import qn
        import re

        # 先清空段落已有内容
        for run in p.runs:
            run.text = ""

        parts = re.split(r'(\-?\d+(?:[,，]\d+)*(?:\.\d+)?(?:[%％])?)', text)
        first_run = True
        for part in parts:
            if not part:
                continue
            is_num = bool(re.match(r'^-?\d+(?:[,，]\d+)*(?:\.\d+)?(?:[%％])?$', part))
            if first_run:
                run = p.add_run()
                first_run = False
            else:
                run = p.add_run()
            run.text = part
            run.font.name = font_name
            run.font.color.rgb = color
            if is_num:
                run.font.size = Pt(base_font_size + 3)
                run.font.bold = True
            else:
                run.font.size = Pt(base_font_size)
                run.font.bold = False

    def _add_bullet_list(self, slide, left, top, width, height, items: List[str],
                         font_size=15, color=None, icon="▸"):
        """添加要点列表"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            full_text = f"{icon} {item}"
            # 清空默认 text 并用带数字强调的方式填充
            p.text = ""
            self._set_paragraph_with_numeric_bold(
                p, full_text, font_size, self.FONT_BODY, color or self.COLOR_LIGHT
            )
            p.space_after = Pt(9)
        return txBox

    def _add_accent_line(self, slide, left, top, width):
        """添加金色装饰线"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLOR_ACCENT
        shape.line.fill.background()
        return shape

    # ------------------------------------------------------------------
    # 图表辅助方法
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_numeric(value_str: str) -> Optional[float]:
        """从数值字符串中提取第一个数字（含小数/负数）。"""
        m = re.search(r"-?\d+(?:\.\d+)?", value_str.replace(",", ""))
        return float(m.group()) if m else None

    def _should_show_chart(self, category_id: str, data_points: List) -> bool:
        """
        智能判断是否应该显示图表。

        Returns:
            bool - 是否显示图表
        """
        # 强数值化类别（必选图表）
        strong_numeric_cats = {"climate", "hydrology", "geophysics", "remote_sensing"}

        # 不适合图表的类别
        no_chart_cats = {"mining_rights"}

        if category_id in no_chart_cats:
            return False

        if category_id in strong_numeric_cats:
            return True

        # 中等数值化类别：50%以上数据点是数值才显示图表
        numeric_count = sum(1 for dp in data_points[:8]
                           if self._extract_numeric(dp.value) is not None)
        return numeric_count >= len(data_points[:8]) * 0.5

    def _select_chart_type(self, category_id: str, data_points: List) -> str:
        """
        根据类别和数据特征选择图表类型。

        Returns:
            str - 图表类型：'bar', 'pie', 'line', 'flow'
        """
        # 检查是否有时间/季节数据
        has_time_data = any(
            '月' in dp.value or '季' in dp.value or '年' in dp.value or
            '春' in dp.value or '夏' in dp.value or '秋' in dp.value or '冬' in dp.value
            for dp in data_points
        )

        # 检查是否有占比/百分比数据
        has_percentage_data = any(
            '%' in dp.value or '占比' in dp.value or '比例' in dp.value
            for dp in data_points
        )

        if category_id == "climate":
            # 气候类：有时间数据用折线图，用饼图展示占比
            return "pie" if has_percentage_data else "line"
        elif category_id == "hydrology":
            # 水文类：流量用过程线，其他用条形图
            return "flow" if any("流量" in dp.item or "径流" in dp.item for dp in data_points) else "bar"
        elif category_id == "geography":
            # 地理类：占比数据用饼图，其他用条形图
            return "pie" if has_percentage_data else "bar"
        elif category_id == "remote_sensing":
            # 遥感类：NDVI/时序用折线图，分类占比用饼图
            return "line" if any("NDVI" in dp.item or "植被" in dp.item for dp in data_points) else "pie"
        else:
            # 其他类别默认用条形图
            return "bar"

    def _group_data_by_dimension(self, data_points: List) -> List[dict]:
        """
        将数据按维度分组（温度、降水、海拔等）。
        返回 [{dimension, label, value}, ...]
        """
        dimension_keywords = {
            "温度": ["气温", "温度", "均温", "最冷", "最热", "高温", "低温"],
            "降水": ["降雨", "降水", "蒸发", "湿度", "雨量"],
            "地形": ["海拔", "高程", "地形", "地势", "高程"],
            "流量": ["流量", "径流", "水位", "流速"],
            "储量": ["储量", "资源量", "储量", "资源", "吨"],
            "面积": ["面积", "范围", "规模"],
            "比例": ["占比", "比例", "百分", "%"],
            "经济": ["GDP", "产值", "收入", "产量"]
        }

        grouped = []
        for dp in data_points:
            value = self._extract_numeric(dp.value)
            if value is None:
                continue

            # 判断数据属于哪个维度
            dimension = "其他"
            for dim, keywords in dimension_keywords.items():
                if any(keyword in dp.item for keyword in keywords):
                    dimension = dim
                    break

            grouped.append({
                "dimension": dimension,
                "label": dp.item,
                "value": value,
                "original_item": dp.item
            })

        return grouped

    def _try_build_chart_data(self, category_id: str, data_points) -> Optional[dict]:
        """
        尝试从 data_points 提取可绘图的数值序列。
        现在会考虑数据维度，只将相同维度的数据放入同一图表。
        返回 {type, labels, values, dimension} 或 None（退回表格）。
        """
        chart_type = self._select_chart_type(category_id, data_points)

        # 先按维度分组
        grouped_data = self._group_data_by_dimension(data_points)
        if not grouped_data:
            return None

        # 检查主要维度
        dimension_count = {}
        for item in grouped_data:
            dim = item["dimension"]
            dimension_count[dim] = dimension_count.get(dim, 0) + 1

        # 如果数据分散在3个以上维度，不适合做单图表
        if len(dimension_count) > 2:
            print(f"[Chart] {category_id} 数据维度过多({len(dimension_count)}个)，跳过图表")
            return None

        # 取数量最多的维度
        main_dimension = max(dimension_count, key=dimension_count.get)
        filtered_data = [item for item in grouped_data if item["dimension"] == main_dimension]

        if len(filtered_data) < 3:
            return None

        labels = [item["label"][:8] if chart_type == "pie" else item["label"][:10]
                 for item in filtered_data[:8]]
        values = [item["value"] for item in filtered_data[:8]]

        return {
            "type": chart_type,
            "labels": labels,
            "values": values,
            "dimension": main_dimension
        }

    def _add_bar_chart(self, slide, left, top, width, height,
                       labels: List[str], values: List[float], title: str = ""):
        """
        在幻灯片上添加横向条形图（深色主题）。
        图表背景透明，轴标签白色，系列用金色渐变填充。
        """
        chart_data = ChartData()
        chart_data.categories = labels
        chart_data.add_series("数值", values)

        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            left, top, width, height,
            chart_data
        )
        chart = chart_shape.chart

        # ---- 绘图区背景 ----
        try:
            chart.plot_area.format.fill.solid()
            chart.plot_area.format.fill.fore_color.rgb = self.COLOR_BG_LIGHT
        except Exception:
            pass

        # ---- 系列颜色：金色 ----
        series = chart.series[0]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_ACCENT

        # ---- 图例隐藏 ----
        chart.has_legend = False

        # ---- 标题 ----
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
            chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.COLOR_LIGHT
            chart.chart_title.text_frame.paragraphs[0].font.name = self.FONT_BODY
        else:
            chart.has_title = False

        # ---- 坐标轴字体白色 ----
        try:
            for axis in (chart.category_axis, chart.value_axis):
                axis.tick_labels.font.size = Pt(9)
                axis.tick_labels.font.color.rgb = self.COLOR_LIGHT
                axis.tick_labels.font.name = self.FONT_BODY
                # 轴线颜色
                axis.format.line.color.rgb = self.COLOR_BG_LIGHT
        except Exception:
            pass

        return chart_shape

    def _add_pie_chart(self, slide, left, top, width, height,
                      labels: List[str], values: List[float], title: str = ""):
        """
        在幻灯片上添加饼图（深色主题）。
        使用金色配色方案。
        """
        chart_data = ChartData()
        chart_data.categories = labels
        chart_data.add_series("占比", values)

        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            left, top, width, height,
            chart_data
        )
        chart = chart_shape.chart

        # ---- 绘图区背景 ----
        try:
            chart.plot_area.format.fill.solid()
            chart.plot_area.format.fill.fore_color.rgb = self.COLOR_BG_LIGHT
        except Exception:
            pass

        # ---- 系列颜色：金色渐变 ----
        series = chart.series[0]
        for i, point in enumerate(series.points):
            fill = point.format.fill
            fill.solid()
            # 使用不同深浅的金色
            gold_shades = [
                RGBColor(0xD4, 0xA8, 0x43),  # 主金色
                RGBColor(0xE5, 0xB7, 0x4C),  # 亮金色
                RGBColor(0xC1, 0x8A, 0x34),  # 暗金色
                RGBColor(0xF4, 0xC8, 0x3B),  # 浅金色
                RGBColor(0xB5, 0x8A, 0x26),  # 深金色
            ]
            fill.fore_color.rgb = gold_shades[i % len(gold_shades)]

        # ---- 图例 ----
        chart.has_legend = True
        if chart.legend:
            chart.legend.position = "r"  # 右侧
            legend = chart.legend.legend_entries
            for entry in legend:
                entry.font.size = Pt(9)
                entry.font.color.rgb = self.COLOR_LIGHT
                entry.font.name = self.FONT_BODY

        # ---- 数据标签 ----
        data_labels = chart.series[0].data_labels
        data_labels.show_percent = True
        data_labels.show_category_name = False
        data_labels.show_legend_key = False
        if data_labels:
            for dl in data_labels:
                dl.font.size = Pt(8)
                dl.font.color.rgb = self.COLOR_WHITE
                dl.font.name = self.FONT_BODY

        # ---- 标题 ----
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
            chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.COLOR_LIGHT
            chart.chart_title.text_frame.paragraphs[0].font.name = self.FONT_BODY
        else:
            chart.has_title = False

        return chart_shape

    def _add_line_chart(self, slide, left, top, width, height,
                       labels: List[str], values: List[float], title: str = ""):
        """
        在幻灯片上添加折线图（深色主题）。
        金色线条，数据点标记。
        """
        chart_data = ChartData()
        chart_data.categories = labels
        chart_data.add_series("数值", values)

        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            left, top, width, height,
            chart_data
        )
        chart = chart_shape.chart

        # ---- 绘图区背景 ----
        try:
            chart.plot_area.format.fill.solid()
            chart.plot_area.format.fill.fore_color.rgb = self.COLOR_BG_LIGHT
        except Exception:
            pass

        # ---- 线条样式 ----
        series = chart.series[0]
        line = series.format.line
        line.color.rgb = self.COLOR_ACCENT
        line.width = Pt(2.5)

        # 数据点样式
        if series.points:
            for point in series.points:
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = self.COLOR_ACCENT
                border = point.format.line
                border.color.rgb = self.COLOR_WHITE
                border.width = Pt(1)

        # ---- 图例隐藏 ----
        chart.has_legend = False

        # ---- 标题 ----
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
            chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.COLOR_LIGHT
            chart.chart_title.text_frame.paragraphs[0].font.name = self.FONT_BODY
        else:
            chart.has_title = False

        # ---- 坐标轴字体白色 ----
        try:
            for axis in (chart.category_axis, chart.value_axis):
                axis.tick_labels.font.size = Pt(9)
                axis.tick_labels.font.color.rgb = self.COLOR_LIGHT
                axis.tick_labels.font.name = self.FONT_BODY
                # 轴线颜色
                axis.format.line.color.rgb = self.COLOR_BG_LIGHT
        except Exception:
            pass

        return chart_shape

    def _add_flow_chart(self, slide, left, top, width, height,
                       labels: List[str], values: List[float], title: str = ""):
        """
        在幻灯片上添加流程线图（模拟流量过程线，深色主题）。
        金色填充，面积图样式。
        """
        chart_data = ChartData()
        chart_data.categories = labels
        chart_data.add_series("流量", values)

        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.AREA,
            left, top, width, height,
            chart_data
        )
        chart = chart_shape.chart

        # ---- 绘图区背景 ----
        try:
            chart.plot_area.format.fill.solid()
            chart.plot_area.format.fill.fore_color.rgb = self.COLOR_BG_LIGHT
        except Exception:
            pass

        # ---- 填充区域：金色渐变 ----
        series = chart.series[0]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_ACCENT
        fill.transparency = 0.3  # 半透明

        # 边框线条
        line = series.format.line
        line.color.rgb = self.COLOR_ACCENT
        line.width = Pt(2)

        # ---- 图例隐藏 ----
        chart.has_legend = False

        # ---- 标题 ----
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
            chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.COLOR_LIGHT
            chart.chart_title.text_frame.paragraphs[0].font.name = self.FONT_BODY
        else:
            chart.has_title = False

        # ---- 坐标轴字体白色 ----
        try:
            for axis in (chart.category_axis, chart.value_axis):
                axis.tick_labels.font.size = Pt(9)
                axis.tick_labels.font.color.rgb = self.COLOR_LIGHT
                axis.tick_labels.font.name = self.FONT_BODY
                # 轴线颜色
                axis.format.line.color.rgb = self.COLOR_BG_LIGHT
        except Exception:
            pass

        return chart_shape

    def _add_table_shape(self, slide, left, top, width, height,
                         headers: List[str], rows: List[List[str]]):
        """添加带样式的数据表格"""
        n_rows = len(rows) + 1
        n_cols = len(headers)
        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = table_shape.table

        # 设置列宽比例（2 列：项目/数值；3 列：项目/数值/其他）
        col_widths = [0.35, 0.65] if n_cols == 2 else [0.3, 0.4, 0.3][:n_cols]
        total = sum(col_widths)
        for i, w in enumerate(col_widths[:n_cols]):
            table.columns[i].width = int(width * w / total)

        # 表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._style_cell(cell, font_size=13, bold=True,
                             font_color=self.COLOR_ACCENT, bg_color=self.COLOR_TABLE_HEADER)

        # 数据行
        for r, row_data in enumerate(rows, start=1):
            bg = self.COLOR_TABLE_ROW1 if r % 2 == 1 else self.COLOR_TABLE_ROW2
            for c, text in enumerate(row_data):
                cell = table.cell(r, c)
                cell.text = text
                self._style_cell(cell, font_size=12, font_color=self.COLOR_LIGHT, bg_color=bg)

        return table_shape

    def _style_cell(self, cell, font_size=10, bold=False, font_color=None, bg_color=None):
        """设置单元格样式"""
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = self.FONT_BODY
            paragraph.font.bold = bold
            paragraph.font.color.rgb = font_color or self.COLOR_WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color

    # ------------------------------------------------------------------
    # 幻灯片页面
    # ------------------------------------------------------------------

    def _add_cover_slide(self, prs, location: LocationContext, mineral_type: str):
        """封面页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        self._set_slide_bg(slide)

        # 顶部装饰线
        self._add_accent_line(slide, Inches(2), Inches(2.2), Inches(9.333))

        # 主标题
        self._add_textbox(
            slide, Inches(2), Inches(2.5), Inches(9.333), Inches(1.2),
            "地质勘探综合报告", font_size=40, font_name=self.FONT_TITLE,
            color=self.COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER
        )

        # 副标题（研究区名称）
        self._add_textbox(
            slide, Inches(2), Inches(3.7), Inches(9.333), Inches(0.8),
            location.area_name, font_size=28, font_name=self.FONT_TITLE,
            color=self.COLOR_ACCENT, alignment=PP_ALIGN.CENTER
        )

        # 信息行
        info_parts = [location.location_str]
        if mineral_type:
            info_parts.append(f"目标矿种：{mineral_type}")
        info_parts.append(datetime.now().strftime("%Y年%m月%d日"))
        info_text = "  |  ".join(info_parts)

        self._add_textbox(
            slide, Inches(2), Inches(4.8), Inches(9.333), Inches(0.6),
            info_text, font_size=14, color=self.COLOR_LIGHT, alignment=PP_ALIGN.CENTER
        )

        # 底部装饰线
        self._add_accent_line(slide, Inches(2), Inches(5.5), Inches(9.333))

    def _generate_location_map(self, location: LocationContext,
                                width_px: int = 660, height_px: int = 600) -> Optional[str]:
        """
        拼接高德地图瓦片（style=8 中文标注），绘制红框标注研究区，
        返回临时 PNG 文件路径；失败返回 None。
        """
        import math
        import urllib.request
        import ssl

        try:
            from PIL import Image, ImageDraw
        except ImportError:
            return None

        try:
            min_lon = location.min_lon
            min_lat = location.min_lat
            max_lon = location.max_lon
            max_lat = location.max_lat
            center_lat = location.centroid_lat
            center_lon = location.centroid_lon

            # 根据 bbox 跨度选缩放级别
            span = max(max_lon - min_lon, max_lat - min_lat)
            if span < 0.03:
                zoom = 14
            elif span < 0.1:
                zoom = 12
            elif span < 0.5:
                zoom = 10
            elif span < 2.0:
                zoom = 8
            elif span < 8.0:
                zoom = 7
            else:
                zoom = 6

            TILE_SIZE = 256

            def lon_to_tile_x(lon):
                return (lon + 180.0) / 360.0 * (2 ** zoom)

            def lat_to_tile_y(lat):
                lat_r = math.radians(lat)
                return (1.0 - math.log(math.tan(lat_r) + 1.0 / math.cos(lat_r)) / math.pi) / 2.0 * (2 ** zoom)

            # 中心瓦片（浮点）
            cx_f = lon_to_tile_x(center_lon)
            cy_f = lat_to_tile_y(center_lat)

            # 需要多少列/行瓦片
            cols = math.ceil(width_px / TILE_SIZE) + 2
            rows = math.ceil(height_px / TILE_SIZE) + 2

            # 左上角起始瓦片
            start_x = int(cx_f - cols / 2)
            start_y = int(cy_f - rows / 2)

            # 拼接画布（比目标稍大，之后裁剪）
            canvas_w = cols * TILE_SIZE
            canvas_h = rows * TILE_SIZE
            canvas = Image.new("RGB", (canvas_w, canvas_h), (200, 200, 200))

            ssl_ctx = ssl.create_default_context()
            ssl_ctx.check_hostname = False
            ssl_ctx.verify_mode = ssl.CERT_NONE

            max_tile = 2 ** zoom
            for tx in range(cols):
                for ty in range(rows):
                    tile_x = (start_x + tx) % max_tile
                    tile_y = (start_y + ty) % max_tile
                    if tile_y < 0 or tile_y >= max_tile:
                        continue
                    # 高德卫星地形图（style=6）
                    url_sat = (f"https://webst01.is.autonavi.com/appmaptile"
                               f"?style=6&x={tile_x}&y={tile_y}&z={zoom}")
                    # 叠加中文标注层（style=8，透明底）
                    url_label = (f"https://wprd01.is.autonavi.com/appmaptile"
                                 f"?lang=zh_cn&size=1&scl=2&style=8"
                                 f"&x={tile_x}&y={tile_y}&z={zoom}")
                    try:
                        req = urllib.request.Request(url_sat, headers={"User-Agent": "geo-reporter/0.1"})
                        with urllib.request.urlopen(req, timeout=8, context=ssl_ctx) as resp:
                            tile_img = Image.open(io.BytesIO(resp.read())).convert("RGBA")
                        try:
                            req2 = urllib.request.Request(url_label, headers={"User-Agent": "geo-reporter/0.1"})
                            with urllib.request.urlopen(req2, timeout=8, context=ssl_ctx) as resp2:
                                label_img = Image.open(io.BytesIO(resp2.read())).convert("RGBA")
                            tile_img = Image.alpha_composite(tile_img, label_img)
                        except Exception:
                            pass  # 标注失败仍用卫星底图
                        canvas.paste(tile_img.convert("RGB"), (tx * TILE_SIZE, ty * TILE_SIZE))
                    except Exception:
                        pass  # 个别瓦片失败保留灰色

            # 计算中心像素在画布上的位置
            center_canvas_x = (cx_f - start_x) * TILE_SIZE
            center_canvas_y = (cy_f - start_y) * TILE_SIZE

            # 裁剪到目标尺寸（以中心为基准）
            crop_x0 = int(center_canvas_x - width_px / 2)
            crop_y0 = int(center_canvas_y - height_px / 2)
            crop_x1 = crop_x0 + width_px
            crop_y1 = crop_y0 + height_px
            img = canvas.crop((crop_x0, crop_y0, crop_x1, crop_y1))

            # ── 坐标转像素（相对于裁剪后图像） ──
            def geo_to_px(lon, lat):
                fx = (lon_to_tile_x(lon) - start_x) * TILE_SIZE - crop_x0
                fy = (lat_to_tile_y(lat) - start_y) * TILE_SIZE - crop_y0
                return fx, fy

            draw = ImageDraw.Draw(img)

            bx0, by0 = geo_to_px(min_lon, max_lat)  # 左上
            bx1, by1 = geo_to_px(max_lon, min_lat)  # 右下

            # 保证最小可见尺寸
            if abs(bx1 - bx0) < 20:
                mx = (bx0 + bx1) / 2; bx0, bx1 = mx - 10, mx + 10
            if abs(by1 - by0) < 20:
                my = (by0 + by1) / 2; by0, by1 = my - 10, my + 10

            # 红框（细线，2px）
            draw.rectangle([bx0, by0, bx1, by1], outline=(220, 30, 30), width=2)


            # 金色外框（直接画在图上，兼容 WPS）
            border_lw = 3
            for offset in range(border_lw):
                draw.rectangle([offset, offset, width_px - 1 - offset, height_px - 1 - offset],
                                outline=(212, 168, 67))

            tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            img.save(tmp.name, "PNG")
            tmp.close()
            return tmp.name

        except Exception:
            return None

    def _add_overview_slide(self, prs, location: LocationContext, mineral_type: str):
        """研究区概况页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide)

        # 标题
        self._add_textbox(
            slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
            "研究区概况", font_size=28, font_name=self.FONT_TITLE,
            color=self.COLOR_ACCENT, bold=True
        )
        self._add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3))

        # 左侧信息卡片
        info_items = [
            f"位置：{location.location_str}",
            f"国家：{location.country}",
            f"省份：{location.province or '未确定'}",
            f"城市：{location.city or '未确定'}",
            f"区县：{location.district or '未确定'}",
            f"坐标范围：{location.coords_str}",
        ]
        if mineral_type:
            info_items.append(f"目标矿种：{mineral_type}")

        self._add_bullet_list(
            slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5),
            info_items, font_size=17, icon="●"
        )

        # 右侧：尝试生成地图，失败则显示占位框
        map_left = Inches(7)
        map_top = Inches(1.5)
        map_width = Inches(5.5)
        map_height = Inches(5)

        map_path = self._generate_location_map(location, width_px=660, height_px=600)
        if map_path:
            try:
                slide.shapes.add_picture(map_path, map_left, map_top, map_width, map_height)
                # 图题
                self._add_textbox(
                    slide, map_left, map_top + map_height, map_width, Inches(0.3),
                    f"图：{location.area_name} 研究区位置示意图",
                    font_size=10, font_name=self.FONT_BODY,
                    color=self.COLOR_LIGHT, bold=False
                )
            except Exception:
                map_path = None

        if not map_path:
            # 降级：显示占位符
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, map_left, map_top, map_width, map_height
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = self.COLOR_BG_LIGHT
            placeholder.line.color.rgb = self.COLOR_ACCENT
            placeholder.line.width = Pt(1)
            tf = placeholder.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "[ 研究区位置示意图 ]"
            p.font.size = Pt(14)
            p.font.name = self.FONT_TITLE
            p.font.color.rgb = self.COLOR_LIGHT
            p.alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].space_before = Pt(80)

    def _add_impact_box(self, slide, left, top, width, height, mineral_type: str, impact_text: str):
        """添加矿种影响专用高亮框（金色边框 + 稍浅背景 + 左侧彩条强调）"""
        # 左侧金色竖条装饰
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Pt(6), height
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.COLOR_ACCENT
        bar.line.fill.background()

        # 主框体（留出左侧竖条位置）
        box_left = left + Pt(10)
        box_width = width - Pt(10)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, box_left, top, box_width, height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x2A, 0x1F, 0x05)  # 深金棕底色
        box.line.color.rgb = self.COLOR_ACCENT
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(12)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(5)
        tf.margin_bottom = Pt(5)

        # 标签行（加大字号）
        p_label = tf.paragraphs[0]
        p_label.text = f"⚠ 对 {mineral_type} 勘探作业的影响"
        p_label.font.size = Pt(14)
        p_label.font.name = self.FONT_TITLE
        p_label.font.color.rgb = self.COLOR_ACCENT
        p_label.font.bold = True
        p_label.space_after = Pt(5)

        # 内容行（加大字号 + 数字强调）
        p_body = tf.add_paragraph()
        p_body.text = ""
        self._set_paragraph_with_numeric_bold(
            p_body, impact_text, 13, self.FONT_BODY, self.COLOR_WHITE
        )

    def _add_category_slide(self, prs, title: str, result: SearchResult, mineral_type: str, category_id: str):
        """数据类别页：概述 + 左表格 + 右要点 + 矿种影响"""
        # 气候页使用专属布局
        if category_id == "climate":
            self._add_climate_slide(prs, title, result, mineral_type)
            return

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide)

        # 标题
        self._add_textbox(
            slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.65),
            title, font_size=24, font_name=self.FONT_TITLE,
            color=self.COLOR_ACCENT, bold=True
        )
        self._add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(11.5))

        # ---- 概述文本（summary）----
        y_cursor = Inches(1.05)
        if result.summary:
            # 提炼梗概（优先含数值的句子，不超过120字）
            summary_text = self._condense_summary(result.summary)
            self._add_textbox(
                slide, Inches(0.8), y_cursor, Inches(11.5), Inches(0.85),
                summary_text, font_size=13, color=self.COLOR_LIGHT
            )
            y_cursor += Inches(0.9)

        # ---- 中部区域：左表格 + 右要点 ----
        has_table = bool(_valid_dps(result))
        has_findings = bool(result.key_findings)
        has_impact = bool(result.exploration_impact and mineral_type)

        # 底部留给矿种影响框的高度
        impact_height = Inches(1.35) if has_impact else Inches(0)
        bottom_margin = Inches(0.25)
        mid_bottom = self.SLIDE_HEIGHT - impact_height - bottom_margin - (Inches(0.15) if has_impact else 0)
        mid_height = mid_bottom - y_cursor

        if has_table and has_findings:
            # 智能判断是否应该显示图表
            should_chart = self._should_show_chart(category_id, result.data_points)

            if should_chart:
                # 提取图表数据并选择类型
                chart_info = self._try_build_chart_data(category_id, result.data_points)
                if chart_info:
                    # 计算图表尺寸（饼图需要更大空间）
                    chart_type = chart_info["type"]
                    if chart_type == "pie":
                        # 饼图需要更多宽度
                        chart_width = Inches(6.0)
                        chart_height = Inches(5.0)
                        left_pos = Inches(0.3)
                    else:
                        # 其他类型用标准尺寸
                        chart_width = Inches(6.8)
                        chart_height = min(Inches(5.0), mid_height)
                        left_pos = Inches(0.4)

                    # 根据类型调用相应图表函数
                    if chart_type == "bar":
                        chart_func = self._add_bar_chart
                    elif chart_type == "pie":
                        chart_func = self._add_pie_chart
                    elif chart_type == "line":
                        chart_func = self._add_line_chart
                    elif chart_type == "flow":
                        chart_func = self._add_flow_chart
                    else:
                        chart_func = self._add_bar_chart

                    # 添加图表
                    dimension_name = chart_info.get("dimension", "")
                    chart_title = f"{dimension_name}分布" if dimension_name else f"{title}分布"

                    chart_func(
                        slide, left_pos, y_cursor, chart_width, chart_height,
                        chart_info["labels"], chart_info["values"],
                        chart_title
                    )

                    # 图表说明文字（仅饼图添加）
                    if chart_type == "pie":
                        dimension_name = chart_info.get("dimension", "")
                        desc_text = f"饼图展示了{dimension_name}的各类别占比分布" if dimension_name else f"饼图展示了{title}的各类别占比分布"
                        self._add_textbox(
                            slide, left_pos, y_cursor + chart_height + Inches(0.2),
                            Inches(6.0), Inches(0.5),
                            desc_text, font_size=10, color=self.COLOR_LIGHT,
                            alignment=PP_ALIGN.CENTER
                        )
                        y_cursor += Inches(0.7)

                    # 右侧要点列表
                    right_pos = Inches(7.8) if chart_type != "pie" else Inches(7.0)
                    right_width = Inches(5.0) if chart_type != "pie" else Inches(5.5)

                    self._add_bullet_list(
                        slide, right_pos, y_cursor, right_width, mid_height,
                        result.key_findings[:6], font_size=15
                    )
                else:
                    # 图表生成失败，回退到表格
                    table_rows = [
                        [dp.item, dp.value]
                        for dp in _valid_dps(result)[:7]
                    ]
                    tbl_height = min(Inches(0.4 + 0.38 * len(table_rows)), mid_height)
                    self._add_table_shape(
                        slide, Inches(0.5), y_cursor, Inches(6.8), tbl_height,
                        ["项目", "数值/描述"], table_rows
                    )
                    self._add_bullet_list(
                        slide, Inches(7.6), y_cursor, Inches(5.3), mid_height,
                        result.key_findings[:6], font_size=15
                    )
            else:
                # 不适合图表：左半表格，右半要点
                table_rows = [
                    [dp.item, dp.value]
                    for dp in _valid_dps(result)[:7]
                ]
                tbl_height = min(Inches(0.4 + 0.38 * len(table_rows)), mid_height)
                self._add_table_shape(
                    slide, Inches(0.5), y_cursor, Inches(6.8), tbl_height,
                    ["项目", "数值/描述"], table_rows
                )
                self._add_bullet_list(
                    slide, Inches(7.6), y_cursor, Inches(5.3), mid_height,
                    result.key_findings[:6], font_size=15
                )

        elif has_table:
            table_rows = [
                [dp.item, dp.value]
                for dp in _valid_dps(result)[:8]
            ]
            row_count = len(table_rows)
            tbl_height = min(Inches(0.4 + 0.38 * row_count), mid_height)
            self._add_table_shape(
                slide, Inches(0.5), y_cursor, Inches(12.3), tbl_height,
                ["项目", "数值/描述"], table_rows
            )

        elif has_findings:
            self._add_bullet_list(
                slide, Inches(0.8), y_cursor, Inches(11.5), mid_height,
                result.key_findings[:7], font_size=16
            )

        # ---- 矿种影响框 ----
        if has_impact:
            impact_top = self.SLIDE_HEIGHT - impact_height - bottom_margin
            self._add_impact_box(
                slide, Inches(0.5), impact_top, Inches(12.3), impact_height,
                mineral_type, result.exploration_impact
            )

    def _add_climate_slide(self, prs, title: str, result: SearchResult, mineral_type: str):
        """气候页专属布局：梗概 + 温度直方图 + 文字气候指标 + 要点 + 矿种影响"""
        import re

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide)

        # 标题
        self._add_textbox(
            slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.65),
            title, font_size=24, font_name=self.FONT_TITLE,
            color=self.COLOR_ACCENT, bold=True
        )
        self._add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(11.5))

        # ---- 梗概（提炼 summary 关键数值句，非平铺全文）----
        y_cursor = Inches(1.05)
        if result.summary:
            summary_text = self._condense_summary(result.summary)
            self._add_textbox(
                slide, Inches(0.8), y_cursor, Inches(11.5), Inches(0.85),
                summary_text, font_size=13, color=self.COLOR_LIGHT
            )
            y_cursor += Inches(0.9)

        has_impact = bool(result.exploration_impact and mineral_type)
        impact_height = Inches(1.35) if has_impact else Inches(0)
        bottom_margin = Inches(0.25)
        mid_bottom = self.SLIDE_HEIGHT - impact_height - bottom_margin - (Inches(0.15) if has_impact else 0)
        mid_height = mid_bottom - y_cursor

        # ---- 从 data_points 中分拣温度 vs 其他气候指标 ----
        TEMP_KEYWORDS = ["极端最高", "极端最低", "7月均温", "7月平均气温", "一月", "1月", "年平均气温", "年均气温", "平均气温"]
        OTHER_KEYWORDS = ["降水天数", "降水量", "日照", "蒸发", "湿度", "风速", "无霜期", "冻土", "7月均降水", "降雨"]

        temp_points = []
        other_points = []
        for dp in result.data_points:
            item_lower = dp.item
            is_temp = any(k in item_lower for k in TEMP_KEYWORDS)
            is_other = any(k in item_lower for k in OTHER_KEYWORDS)
            if is_temp:
                temp_points.append(dp)
            elif is_other:
                other_points.append(dp)
            else:
                # 剩余含气温/温度关键词的也归温度
                if "气温" in item_lower or "温度" in item_lower:
                    temp_points.append(dp)
                else:
                    other_points.append(dp)

        # 如果筛选后温度点不足，放宽：把所有数值型 data_points 都视为温度
        if len(temp_points) < 2:
            temp_points = result.data_points
            other_points = []

        # ---- 提取温度数值 ----
        def extract_number(val_str: str) -> Optional[float]:
            m = re.search(r'-?\d+(?:\.\d+)?', val_str.replace(',', ''))
            return float(m.group()) if m else None

        chart_labels, chart_values = [], []
        for dp in temp_points[:5]:
            val = extract_number(dp.value)
            if val is not None:
                # 标签最多6字
                chart_labels.append(dp.item[:6])
                chart_values.append(val)

        # ---- 左侧：温度直方图 ----
        chart_width = Inches(6.5)
        chart_height = min(Inches(4.2), mid_height)
        if len(chart_values) >= 2:
            self._add_bar_chart(
                slide, Inches(0.4), y_cursor, chart_width, chart_height,
                chart_labels, chart_values, "气温指标（℃）"
            )
        else:
            # 无数值则显示提示
            self._add_textbox(
                slide, Inches(0.4), y_cursor, chart_width, chart_height,
                "（气温数值数据不足，无法绘图）", font_size=13, color=self.COLOR_LIGHT
            )

        # ---- 右侧上：其他气候指标（文字） ----
        right_x = Inches(7.3)
        right_w = Inches(5.5)
        right_y = y_cursor

        if other_points:
            text_items = [f"{dp.item}：{dp.value}" for dp in other_points[:5]]
            self._add_bullet_list(
                slide, right_x, right_y, right_w, Inches(2.2),
                text_items, font_size=14, icon="◆"
            )
            right_y += Inches(2.3)

        # ---- 右侧下：关键发现 ----
        findings_height = mid_bottom - right_y
        if result.key_findings and findings_height > Inches(0.5):
            self._add_bullet_list(
                slide, right_x, right_y, right_w, findings_height,
                result.key_findings[:4], font_size=14
            )

        # ---- 矿种影响框 ----
        if has_impact:
            impact_top = self.SLIDE_HEIGHT - impact_height - bottom_margin
            self._add_impact_box(
                slide, Inches(0.5), impact_top, Inches(12.3), impact_height,
                mineral_type, result.exploration_impact
            )

    def _condense_summary(self, summary: str) -> str:
        """
        将 summary 提炼为梗概：
        - 优先取含数值的句子（最多2句）
        - 总长度控制在120字以内
        """
        import re
        sentences = re.split(r'[。；;]', summary)
        # 优先选含数字的句子
        numeric_sents = [s.strip() for s in sentences if re.search(r'\d', s) and len(s.strip()) > 5]
        other_sents = [s.strip() for s in sentences if not re.search(r'\d', s) and len(s.strip()) > 5]

        result_sents = numeric_sents[:2]
        if len(result_sents) < 2:
            result_sents += other_sents[:2 - len(result_sents)]

        condensed = "；".join(result_sents[:2])
        if len(condensed) > 120:
            condensed = condensed[:119] + "…"
        return condensed or summary[:120]

    def _fig_attr(self, fig, key):
        return getattr(fig, key, None) if not isinstance(fig, dict) else fig.get(key)

    def _add_figures_slide(self, prs, title: str, figures: list):
        """子系统图件页：最多 2 张图并排，附图注。"""
        import os
        figs = [f for f in figures if self._fig_attr(f, "path")
                and os.path.exists(self._fig_attr(f, "path"))]
        if not figs:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide)
        self._add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.7),
                          f"{title} · 图件", font_size=26, font_name=self.FONT_TITLE,
                          color=self.COLOR_ACCENT, bold=True)
        self._add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3))

        shown = figs[:2]
        img_w = Inches(5.6)
        gap = Inches(0.4)
        total_w = img_w * len(shown) + gap * (len(shown) - 1)
        start_left = (self.SLIDE_WIDTH - total_w) / 2
        for i, f in enumerate(shown):
            left = start_left + i * (img_w + gap)
            try:
                pic = slide.shapes.add_picture(self._fig_attr(f, "path"), left, Inches(1.5), width=img_w)
            except Exception:
                continue
            cap = self._fig_attr(f, 'caption') or ''
            self._add_textbox(slide, left, Inches(6.4), img_w, Inches(0.8), cap,
                              font_size=11, color=self.COLOR_LIGHT, alignment=PP_ALIGN.CENTER)

    def _add_target_slide(self, prs, location, target_figure, mineral_type: str):
        """靶区推荐页：底图框定靶区。"""
        import os
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide)
        self._add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.7),
                          "靶区推荐", font_size=28, font_name=self.FONT_TITLE,
                          color=self.COLOR_ACCENT, bold=True)
        self._add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3))

        path = self._fig_attr(target_figure, "path") if target_figure is not None else None
        targets = getattr(target_figure, "targets", None) if target_figure is not None else None
        if path and os.path.exists(path):
            try:
                # 左侧热力靶区图
                slide.shapes.add_picture(path, Inches(0.6), Inches(1.4), height=Inches(5.2))
            except Exception:
                path = None
            cap = self._fig_attr(target_figure, "caption") or ""
            self._add_textbox(slide, Inches(0.6), Inches(6.7), Inches(7), Inches(0.6),
                              cap, font_size=11, color=self.COLOR_LIGHT, alignment=PP_ALIGN.CENTER)
            # 右侧：各靶区置信评级 + 理由
            if targets:
                grade_color = {"A": "🟥", "B": "🟧", "C": "🟦", "D": "🟦"}
                items = []
                for t in targets:
                    g = t.get("grade", "")
                    items.append(f"靶区#{t.get('rank','')} [{g}级] {t.get('reason','')}")
                self._add_textbox(slide, Inches(8.0), Inches(1.4), Inches(4.9), Inches(0.5),
                                  "靶区置信评级（A>B>C>D）", font_size=15,
                                  color=self.COLOR_ACCENT, bold=True)
                self._add_bullet_list(slide, Inches(8.0), Inches(2.0), Inches(4.9), Inches(4.8),
                                      items, font_size=11)
        if not path or not os.path.exists(path):
            self._add_textbox(slide, Inches(0.8), Inches(3), Inches(11.5), Inches(1),
                              "靶区推荐图生成失败（底图获取异常），请检查网络后重试。",
                              font_size=16, color=self.COLOR_LIGHT, alignment=PP_ALIGN.CENTER)

    def _add_confidence_slide(self, prs, confidence: dict, mineral_type: str):
        """综合置信评价页（A-B-C-D）。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide)
        self._add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.7),
                          "综合置信评价", font_size=28, font_name=self.FONT_TITLE,
                          color=self.COLOR_ACCENT, bold=True)
        self._add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3))

        if not confidence:
            self._add_textbox(slide, Inches(0.8), Inches(3), Inches(11.5), Inches(1),
                              "置信评价未生成。", font_size=16, color=self.COLOR_LIGHT,
                              alignment=PP_ALIGN.CENTER)
            return

        grade = str(confidence.get("grade", "")).strip().upper()[:1]
        label = confidence.get("grade_label", "")
        grade_colors = {"A": RGBColor(0x2E, 0xCC, 0x71), "B": RGBColor(0x34, 0x98, 0xDB),
                        "C": RGBColor(0xF3, 0x9C, 0x12), "D": RGBColor(0xE7, 0x4C, 0x3C)}
        self._add_textbox(slide, Inches(0.8), Inches(1.4), Inches(4), Inches(1.4),
                          grade or "—", font_size=72, font_name=self.FONT_TITLE,
                          color=grade_colors.get(grade, self.COLOR_WHITE), bold=True,
                          alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, Inches(0.8), Inches(2.9), Inches(4), Inches(0.6),
                          label, font_size=18, color=self.COLOR_ACCENT, alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, Inches(0.8), Inches(3.6), Inches(4), Inches(0.5),
                          "A>B>C>D（A 最高）", font_size=12, color=self.COLOR_LIGHT,
                          alignment=PP_ALIGN.CENTER)

        # 右侧：综述 + 分项
        summary = confidence.get("summary", "")
        self._add_textbox(slide, Inches(5.2), Inches(1.4), Inches(7.3), Inches(2.2),
                          summary, font_size=13, color=self.COLOR_WHITE)
        dims = confidence.get("dimensions", [])
        if dims:
            items = [f"{d.get('name','')}：{d.get('level','')} — {d.get('note','')}" for d in dims]
            self._add_bullet_list(slide, Inches(5.2), Inches(3.7), Inches(7.3), Inches(3.2),
                                  items, font_size=12)
        rec = confidence.get("recommendation", "")
        if rec:
            self._add_textbox(slide, Inches(0.8), Inches(4.5), Inches(4), Inches(2.5),
                              f"建议：{rec}", font_size=13, color=self.COLOR_LIGHT)

    def _add_summary_slide(self, prs, search_results: Dict[str, SearchResult], mineral_type: str):
        """综合发现汇总页：左侧发现列表 + 右侧柱状图（各类别数据点数）"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide)

        self._add_textbox(
            slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
            "综合发现汇总", font_size=28, font_name=self.FONT_TITLE,
            color=self.COLOR_ACCENT, bold=True
        )
        self._add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3))

        # 左侧：各类别首条发现
        all_findings = []
        chart_labels, chart_values = [], []
        for category in get_all_categories():
            result = search_results.get(category.id)
            if result and not result.error:
                if result.key_findings:
                    all_findings.append(f"【{category.name}】{result.key_findings[0]}")
                chart_labels.append(category.name[:5])  # 短标签
                chart_values.append(len(result.data_points) if result.data_points else 0)

        if all_findings:
            self._add_bullet_list(
                slide, Inches(0.8), Inches(1.5), Inches(6.8), Inches(5.5),
                all_findings[:8], font_size=14, icon="◆"
            )

        # 右侧：数据点数量柱状图
        if len(chart_values) >= 3 and any(v > 0 for v in chart_values):
            self._add_bar_chart(
                slide, Inches(7.8), Inches(1.5), Inches(5.2), Inches(5.5),
                chart_labels, chart_values
            )

    def _add_sources_slide(self, prs, search_results: Dict[str, SearchResult]):
        """数据来源页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide)

        self._add_textbox(
            slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
            "数据来源", font_size=28, font_name=self.FONT_TITLE,
            color=self.COLOR_ACCENT, bold=True
        )
        self._add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3))

        all_sources = set()
        for result in search_results.values():
            if not result.error and result.data_sources:
                all_sources.update(result.data_sources)

        sources = sorted(all_sources)[:15]  # 限制数量
        if sources:
            numbered = [f"{i+1}. {s}" for i, s in enumerate(sources)]
            self._add_bullet_list(
                slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                numbered, font_size=13, icon=""
            )
        else:
            self._add_textbox(
                slide, Inches(0.8), Inches(2), Inches(11), Inches(1),
                "报告基于 AI 知识库和网络搜索综合整理。",
                font_size=15, color=self.COLOR_LIGHT
            )
